from pptx import Presentation

try:
    prs = Presentation('ScoutAnt_Synopsis.pptx')
    with open('extracted_ppt_text.txt', 'w', encoding='utf-8') as f:
        f.write(f"Total slides: {len(prs.slides)}\n")
        for idx, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            f.write(f"Slide {idx} text: {' | '.join(text).replace(chr(11), ' ')}\n")
except Exception as e:
    print(f"Error: {e}")
