import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

def add_bullet_points(slide, text_list, placeholder_idx=1):
    tf = slide.placeholders[placeholder_idx].text_frame
    tf.text = text_list[0]
    for text in text_list[1:]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0

def create_presentation():
    # Load template
    template_path = r'c:\Users\dhruv\OneDrive\Desktop\ScoutAnt-Ultimate-Analyzer-and-Builder\ScoutAnt_Synopsis_Updated.pptx'
    prs = Presentation(template_path)
    
    # Remove existing slides to start fresh with the theme
    # Note: Deleting slides in a loop requires careful indexing or copying
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ScoutAnt: Ultimate Valorant Analyzer and Team Builder"
    subtitle.text = "Project Synopsis Presentation\nPresented by: [Your Name]\n" + datetime.datetime.now().strftime("%B %Y")

    # --- 1. Introduction ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    content = [
        "ScoutAnt is an intelligent, data-driven system for competitive Valorant analysis.",
        "Automates collection of match data from VLR.gg.",
        "Analyzes player performance (ACS, K/D, Agents, Maps).",
        "Provides actionable strategic intelligence for pre-match planning.",
        "Bridges the gap between raw data and tactical decision-making."
    ]
    add_bullet_points(slide, content)

    # --- 2. Literature Survey / Existing System ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Literature Survey / Existing System"
    content = [
        "Existing platforms (Tracker.gg, VLR.gg) provide extensive raw data.",
        "Current systems lack automated team synergy analysis and predictive modeling.",
        "Manual analysis is time-intensive and prone to human error.",
        "No centralized tool for automated 'Meta Composition' identification.",
        "ScoutAnt fills this gap by providing high-level processed insights."
    ]
    add_bullet_points(slide, content)

    # --- 3. Requirements and Analysis ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Requirements and Analysis"
    content = [
        "Functional: Event discovery, match scraping, data persistence, analysis algorithms.",
        "Non-Functional: Performance (100+ matches/hr), reliability, modularity.",
        "Data Validation: Ensuring completeness (10 players/map) and accuracy.",
        "User Interface: CLI-based with real-time feedback and summary stats."
    ]
    add_bullet_points(slide, content)

    # --- 4. Proposed Methodology ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Proposed Methodology"
    content = [
        "Step 1: Automated Web Scraping using BeautifulSoup and Requests.",
        "Step 2: Hierarchical Data Storage in structured JSON format.",
        "Step 3: Statistical Processing to identify Meta Compositions.",
        "Step 4: Player Performance Prediction based on historical player-agent-map data.",
        "Step 5: Actionable Intelligence Delivery for pre-match planning."
    ]
    add_bullet_points(slide, content)

    # --- 5. Hardware / Software Requirements ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Hardware / Software Requirements"
    content = [
        "Software:",
        "  - Python 3.x, Requests, BeautifulSoup4, python-pptx",
        "  - Windows/Linux OS",
        "Hardware:",
        "  - Standard CPU (i5 or equivalent)",
        "  - 8GB RAM minimum",
        "  - Stable Internet Connection for scraping"
    ]
    add_bullet_points(slide, content)

    # --- 6. Results ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Results"
    content = [
        "Successful extraction of 1000+ match records into JSON.",
        "Accurate identification of top-tier agent compositions per map.",
        "Reliable player performance forecasting using historical stats.",
        "Modular codebase allows easy addition of new scrapers/analyzers.",
        "Generated clean, data-backed strategic reports."
    ]
    add_bullet_points(slide, content)

    # --- 7. Conclusion and Future Work ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion and Future Work"
    content = [
        "Conclusion: ScoutAnt provides a significant competitive edge through data.",
        "Future Work:",
        "  - ML-based Win Probability Predictor.",
        "  - Interactive Web Dashboard (React/Plotly).",
        "  - Real-time tournament analysis and API integration.",
        "  - Automated VOD-timestamp linking."
    ]
    add_bullet_points(slide, content)

    # --- 8. References (APA Style) ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "References (APA Style)"
    content = [
        "VLR.gg. (2024). Valorant Match Statistics. https://www.vlr.gg/",
        "Python Software Foundation. (2024). Python 3 Documentation. https://docs.python.org/3/",
        "Richardson, L. (2024). BeautifulSoup Documentation. https://www.crummy.com/software/BeautifulSoup/",
        "Riot Games. (2024). Valorant Esports. https://valorantesports.com/"
    ]
    add_bullet_points(slide, content)

    # --- 9. Live Demonstration ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Live Demonstration of the Project"
    
    # Text placeholder for demonstration info
    tf = slide.placeholders[1].text_frame
    tf.text = "The following slides show screenshots of the system in action:"
    
    # We will add 3 images in separate slides or as content here.
    # The requirement says "include placeholders for screenshots"
    
    # Screenshot 1: Terminal - Script Execution
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Demo: Script Execution"
    tf = slide1.placeholders[1].text_frame
    tf.text = "[Screenshot Placeholder: Terminal window executing 'python analyzer.py']"
    
    # Screenshot 2: Terminal - Retrieval Output
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Demo: Data Retrieval Output"
    tf = slide2.placeholders[1].text_frame
    tf.text = "[Screenshot Placeholder: Terminal output showing match scraping progress]"
    
    # Screenshot 3: JSON File - Data Structure
    slide3 = prs.slides.add_slide(slide_layout)
    slide3.shapes.title.text = "Demo: Generated JSON Data"
    tf = slide3.placeholders[1].text_frame
    tf.text = "[Screenshot Placeholder: JSON file view (e.g., match_stats_db.json excerpt)]"

    # Save
    output_path = r'c:\Users\dhruv\OneDrive\Desktop\ScoutAnt-Ultimate-Analyzer-and-Builder\ScoutAnt_Synopsis_Final_Updated.pptx'
    prs.save(output_path)
    print(f"✅ Presentation created: {output_path}")

if __name__ == "__main__":
    create_presentation()
