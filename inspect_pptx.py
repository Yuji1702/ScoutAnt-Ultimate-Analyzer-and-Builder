import collections 
import collections.abc
import pptx

prs = pptx.Presentation(r'c:\Users\dhruv\OneDrive\Desktop\ScoutAnt-Ultimate-Analyzer-and-Builder\ScoutAnt_Synopsis_Updated.pptx')

print("Slide Layouts:")
for i, layout in enumerate(prs.slide_layouts):
    print(f"[{i}] {layout.name}")
    for shape in layout.placeholders:
        print(f"  - Placeholder {shape.placeholder_format.idx}: {shape.name}")

print("\nExisting Slides:")
for i, slide in enumerate(prs.slides):
    print(f"Slide {i+1} (Layout: {slide.slide_layout.name})")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(f"  - Text: '{run.text}', Font: {run.font.name}, Size: {run.font.size}, Color: {getattr(run.font.color, 'rgb', 'N/A')}")
        else:
            print(f"  - Shape: {shape.name}")
