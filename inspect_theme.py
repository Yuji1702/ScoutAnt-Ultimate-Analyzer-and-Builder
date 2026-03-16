from pptx import Presentation

prs = Presentation('ScoutAnt_Synopsis.pptx')
slide = prs.slides[2] # 1. Objective slide
print(f"Slide 2")
for idx, shape in enumerate(slide.shapes):
    print(f"Shape {idx}: {shape.name}")
    if hasattr(shape, "text_frame"):
        for p_idx, p in enumerate(shape.text_frame.paragraphs):
            print(f"  Paragraph {p_idx} text: {p.text}")
            if p.runs:
                run = p.runs[0]
                font_name = run.font.name if run.font.name else "Inherited"
                font_size = run.font.size.pt if run.font.size else "Inherited"
                color = "Inherited" 
                try: 
                    if run.font.color and run.font.color.type == 1: # RGB
                        color = run.font.color.rgb
                    elif run.font.color and run.font.color.type == 2: # Theme
                        color = f"Theme {run.font.color.theme_color}"
                except:
                    pass
                print(f"    Font: {font_name}, Size: {font_size}, Color: {color}, Bold: {run.font.bold}")

