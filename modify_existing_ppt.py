from pptx import Presentation
from pptx.util import Pt
import copy

def update_text(shape, new_text, font_name="Alatsi", font_size=20, is_title=False):
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    
    if is_title:
        p = tf.paragraphs[0]
        p.text = new_text
        for run in p.runs:
            run.font.name = "Alatsi"
            run.font.size = Pt(31) if not font_size else Pt(font_size)
    else:
        # Handle bullet points
        lines = new_text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.level = 0
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)

def create_presentation():
    prs = Presentation('ScoutAnt_Synopsis.pptx')
    
    # 1. Objective of the Project -> Introduction
    slide = prs.slides[2]
    update_text(slide.shapes.title, "1. Introduction", is_title=True)
    update_text(slide.shapes[3], "• ScoutAnt is an advanced data-driven system for competitive Valorant analysis.\n" \
                                 "• Automates the collection of comprehensive match data from publicly available sources (VLR.gg).\n" \
                                 "• Aimed at providing insights into meta team compositions, player performance, and match strategy.\n" \
                                 "• Revolutionizes match preparation by removing the need for time-intensive manual analysis.",
                                 font_name="Calibri", font_size=20)
                                 
    # 2. Problem Statement -> Literature Survey / Existing System
    slide = prs.slides[3]
    update_text(slide.shapes.title, "2. Literature Survey / Existing System", is_title=True)
    update_text(slide.shapes[3], "• VLR.gg: Primary source of statistics, but no public API and consumption-only.\n" \
                                 "• Tracker.gg: Focuses heavily on ranked matchmaking, not professional esports.\n" \
                                 "• Rib.gg: Paid service with pre-defined analytical dashboards.\n" \
                                 "• Manual Spreadsheets: Traditional method requiring high effort and time.\n" \
                                 "• ScoutAnt aims to bridge these gaps by offering customizable, automated high-level data pipelines.",
                                 font_name="Calibri", font_size=20)

    # 3. Abstract -> Requirements and Analysis
    slide = prs.slides[4]
    update_text(slide.shapes.title, "3. Requirements and Analysis", is_title=True)
    update_text(slide.shapes[3], "Functional Requirements:\n" \
                                 "  • Event Discovery & Match Extraction (VLR.gg integration)\n" \
                                 "  • Data Persistence (JSON storage system)\n" \
                                 "  • Meta Composition Analysis & Player Prediction\n" \
                                 "Non-Functional Requirements:\n" \
                                 "  • Performance: Rate-limiting to ensure respectful web scraping\n" \
                                 "  • Scalability: Handle thousands of matches efficiently\n" \
                                 "  • Reliability: Graceful exception handling for incomplete data",
                                 font_name="Calibri", font_size=20)

    # 4. Literature Survey -> Proposed Methodology
    slide = prs.slides[5]
    update_text(slide.shapes.title, "4. Proposed Methodology", is_title=True)
    update_text(slide.shapes[3], "• Phase 1: Reconnaissance (Analyzing DOM structure of VLR.gg)\n" \
                                 "• Phase 2: Design Data Architecture (Hierarchical JSON)\n" \
                                 "• Phase 3: Web Scraper Implementation (Python, BeautifulSoup4, Requests)\n" \
                                 "• Phase 4: Data Processing Module (Aggregating ACS, KD, Win Rates)\n" \
                                 "• Phase 5: Result Generation (Identifying Meta Comps)",
                                 font_name="Calibri", font_size=20)

    # 5. Proposed Methodology -> Hardware / Software Requirements
    slide = prs.slides[6]
    update_text(slide.shapes.title, "5. Hardware / Software Requirements", is_title=True)
    update_text(slide.shapes[3], "Hardware Requirements:\n" \
                                 "  • Processor: Multi-core CPU (Intel i3/Ryzen 3) or higher\n" \
                                 "  • RAM: Minimum 4 GB, Recommended 8 GB\n" \
                                 "  • Storage: At least 250 MB free space for JSON storage\n" \
                                 "  • Network: Stable Internet connection for scraping\n" \
                                 "Software Requirements:\n" \
                                 "  • OS: Windows, macOS, or Linux\n" \
                                 "  • Language: Python 3.8+\n" \
                                 "  • Libraries: BeautifulSoup4, Requests",
                                 font_name="Calibri", font_size=16)

    # 6. Requirements -> Results
    slide = prs.slides[7]
    update_text(slide.shapes.title, "6. Results", is_title=True)
    update_text(slide.shapes[3], "• Automated data collection significantly reduces prep time from hours to minutes.\n" \
                                 "• Robust extraction with >99% success rate across standard match pages.\n" \
                                 "• Data successfully structured in 'match_stats_db.json' allowing complex queries.\n" \
                                 "• Achieved 100% accuracy in extracted player stats validated manually.\n" \
                                 "• Predictions show strong correlation with actual tournament outcomes.",
                                 font_name="Calibri", font_size=20)

    # 7. Requirements -> Conclusion and Future Work
    slide = prs.slides[8]
    update_text(slide.shapes.title, "7. Conclusion and Future Work", is_title=True)
    update_text(slide.shapes[3], "Conclusion:\n" \
                                 "  • Democratizes access to high-level esports data through automation.\n" \
                                 "  • Establishes a scalable foundation for Valorant analytics.\n" \
                                 "Future Work:\n" \
                                 "  • Machine Learning integration for win probability prediction.\n" \
                                 "  • Development of a graphical web dashboard or desktop GUI.\n" \
                                 "  • Transition from JSON file storage to a robust relational database.\n" \
                                 "  • Real-time match tracking capabilities.",
                                 font_name="Calibri", font_size=16)

    # 8. Description of Modules -> References
    slide = prs.slides[9]
    update_text(slide.shapes.title, "8. References", is_title=True)
    update_text(slide.shapes[3], "• Python Software Foundation. Python Language Reference. Available at http://www.python.org\n" \
                                 "• Richardson, L. Beautiful Soup Documentation. Retrieved from crummy.com\n" \
                                 "• VLR.gg. Valorant Competitive Match Statistics. Retrieved from https://www.vlr.gg/\n" \
                                 "• Reitz, K. Requests: HTTP for Humans. Retrieved from requests.readthedocs.io",
                                 font_name="Calibri", font_size=16)

    # 9. Result -> Live Demonstration (Adding Screenshots)
    slide = prs.slides[10]
    update_text(slide.shapes.title, "9. Live Demonstration & Screenshots", is_title=True)
    update_text(slide.shapes[3], "", font_name="Calibri") # Clear text for images
    
    from pptx.util import Inches
    try:
        slide.shapes.add_picture('cmd_screenshot.png', Inches(0.5), Inches(2.2), width=Inches(4.5))
        slide.shapes.add_picture('json_screenshot.png', Inches(5.2), Inches(2.2), width=Inches(4.5))
    except Exception as e:
        print(f"Image error: {e}")

    # Delete the remaining unused slides (11, 12, 13) - keeping slide 14 (Thank You)
    # We will just make a list of slide IDs to delete
    # Original slide indices: 0(Title), 1(Index), ..., 10(Live Demo), 11, 12, 13, 14(Thank You)
    
    def safe_delete_slide(prs, slide_index):
        slides = list(prs.slides)
        if slide_index < len(slides):
            slide = slides[slide_index]
            for i, sldId in enumerate(prs.slides._sldIdLst):
                if sldId.id == slide.slide_id:
                    prs.part.drop_rel(sldId.rId)
                    del prs.slides._sldIdLst[i]
                    break
                    
    # Delete backwards to not mess up indices! Delete 13, 12, 11
    safe_delete_slide(prs, 13)
    safe_delete_slide(prs, 12)
    safe_delete_slide(prs, 11)

    # Update Index Slide (Slide 1)
    slide = prs.slides[1]
    update_text(slide.shapes[3], "1. Introduction\n" \
                                 "2. Literature Survey / Existing System\n" \
                                 "3. Requirements and Analysis\n" \
                                 "4. Proposed Methodology\n" \
                                 "5. Hardware / Software Requirements\n" \
                                 "6. Results\n" \
                                 "7. Conclusion and Future Work\n" \
                                 "8. References\n" \
                                 "9. Live Demonstration & Screenshots",
                                 font_name="Calibri", font_size=20)

    prs.save('ScoutAnt_Synopsis_Final_Themed.pptx')
    print("Saved -> ScoutAnt_Synopsis_Final_Themed.pptx")

if __name__ == '__main__':
    create_presentation()
