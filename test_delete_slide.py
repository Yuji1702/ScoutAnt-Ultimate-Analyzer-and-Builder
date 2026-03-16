from pptx import Presentation

def safe_delete_slide(prs, slide_index):
    slides = list(prs.slides)
    if slide_index < len(slides):
        slide = slides[slide_index]
        for i, sldId in enumerate(prs.slides._sldIdLst):
            if sldId.id == slide.slide_id:
                prs.part.drop_rel(sldId.rId)
                del prs.slides._sldIdLst[i]
                break

prs = Presentation('ScoutAnt_Synopsis.pptx')
# delete all but the first slide
for i in range(len(prs.slides) - 1, 0, -1):
    safe_delete_slide(prs, i)

prs.save('test_repair.pptx')
print("Saved test_repair.pptx")
